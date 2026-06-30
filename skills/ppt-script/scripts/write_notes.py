#!/usr/bin/env python3
"""把每页讲稿写进 pptx 备注（speaker notes）。

用法:
    python write_notes.py <deck.pptx> <notes.txt> [-o <out.pptx>] [--in-place]

notes.txt: UTF-8 文本，按幻灯片顺序排列，每页之间用单独一行 `@@@SLIDE@@@` 分隔；
段数必须等于幻灯片数，否则报错退出（不写）。
默认输出 <deck>.noted.pptx（不覆盖定稿）；-o 指定路径；--in-place 原地覆盖。
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation

SENTINEL = "@@@SLIDE@@@"


def split_notes(text: str) -> list[str]:
    """按 SENTINEL 行切分 payload，返回每页备注文本（去首尾空行）。"""
    return [seg.strip("\n") for seg in text.split(SENTINEL)]


def write_notes(deck_path: Path, notes: list[str], out_path: Path) -> int:
    prs = Presentation(str(deck_path))
    n_slides = len(prs.slides)
    if len(notes) != n_slides:
        print(
            f"ERROR: 段数 {len(notes)} != 幻灯片数 {n_slides}；"
            f"请确认每页用单独一行 {SENTINEL} 分隔、段数与页数一致。",
            file=sys.stderr,
        )
        return 1
    for slide, note in zip(prs.slides, notes):
        slide.notes_slide.notes_text_frame.text = note
    prs.save(str(out_path))
    print(f"OK: 已写入 {n_slides} 页备注 -> {out_path}")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description="把讲稿 payload 写进 pptx 备注")
    ap.add_argument("deck", help="输入 pptx 路径")
    ap.add_argument("notes", help="讲稿 payload 文本（@@@SLIDE@@@ 分隔，按页序）")
    ap.add_argument("-o", "--out", help="输出 pptx 路径")
    ap.add_argument("--in-place", action="store_true", help="原地覆盖输入 pptx")
    args = ap.parse_args()

    deck_path = Path(args.deck)
    notes = split_notes(Path(args.notes).read_text(encoding="utf-8"))
    if args.in_place:
        out_path = deck_path
    elif args.out:
        out_path = Path(args.out)
    else:
        out_path = deck_path.with_suffix(".noted.pptx")
    return write_notes(deck_path, notes, out_path)


if __name__ == "__main__":
    sys.exit(main())
