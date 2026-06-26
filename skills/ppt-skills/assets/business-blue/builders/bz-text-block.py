# -*- coding: utf-8 -*-
"""可复用 一段文字 text-block builder（business-blue，蒸自 work_inspire_128 p98/p99）。
版式 = 左「单文本块」(页标题 28 + 大粗引导句 28 + 正文段 16) + 右视觉(手机/图片,可1~2) + 可选放射特性圈(真图标) + 同心 ghost 环 + 底部波浪。
**高频页型**:很多页就是"一个标题 + 一段话 + 配图"。文字/图片属内容、用时传(§6.24);背景复用 bz-light-bg、底波浪复用 bz-wave-layers。
特性圈图标 = 从 ppt-master `templates/icons` 取真图标栅格化成白图标(§资产地图:图标用 templates/icons,别画白圈占位)。

调用：build_text_block(slide, label, lead, body, visuals, features=None, clusters=None)
  visuals=[(x,y,w,h,photo_or_None[,kind,icon]), ...]   features=[(x,y,label,icon_name), ...] 或 None
  clusters=[(x,y,w,h), ...]  ghost 圆簇背景(大圆+周边小圆,整组调库 bz-ghost-cluster;p99 用)
直接运行 = p98 内容 demo → _gallery/text-block.pptx。"""
import os
import sys
import tempfile
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
A1 = "0165FF"
A2 = "00BEFF"
NAVY = "0A2A55"
BODY = "44546A"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
VEC = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "vectors")
ICONLIB = r"D:\Lu.Yao7\github\ppt-master\skills\ppt-master\templates\icons"
ICACHE = os.path.join(tempfile.gettempdir(), "bz_iconcache")
os.makedirs(ICACHE, exist_ok=True)
_parser = etree.XMLParser(resolve_entities=False, no_network=True)
_uid = [5000]


def _nid():
    _uid[0] += 1
    return _uid[0]


def _load(name):
    el = etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    return el


def place(slide, name, x, y, w, h):
    """从素材库载入资产并定位/缩放——**禁内联画图,一律调库**(§资产架构)。
    renumber 全部 cNvPr 避免多实例(如双手机)子件 id 相撞。"""
    el = etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)
    for nv in el.iter(f"{P}cNvPr"):
        nv.set("id", str(_nid()))
    xf = el.find(f"{P}grpSpPr/{A}xfrm") if el.tag == f"{P}grpSp" else el.find(f"{P}spPr/{A}xfrm")
    xf.find(f"{A}off").set("x", str(_e(x))); xf.find(f"{A}off").set("y", str(_e(y)))
    xf.find(f"{A}ext").set("cx", str(_e(w))); xf.find(f"{A}ext").set("cy", str(_e(h)))
    return _ap(slide, el)


def _ap(slide, el):
    slide.shapes._spTree.append(el)
    return el


def _e(v):
    return int(round(v * IN))


def white_icon(name, sub="chunk-filled", scale=8):
    """templates/icons/<sub>/<name>.svg → 白色透明 PNG(缓存),供 add_picture 叠到图标圆。"""
    out = os.path.join(ICACHE, f"{sub}_{name}_w.png")
    if os.path.exists(out):
        return out
    from svglib.svglib import svg2rlg
    from reportlab.graphics import renderPM
    import numpy as np
    from PIL import Image
    d = svg2rlg(os.path.join(ICONLIB, sub, name + ".svg"))
    d.scale(scale, scale); d.width *= scale; d.height *= scale
    raw = out + "_raw.png"
    renderPM.drawToFile(d, raw, fmt="PNG", bg=0xFFFFFF)
    g = np.array(Image.open(raw).convert("L")).astype(np.int16)
    alpha = np.clip(255 - g, 0, 255).astype(np.uint8)
    h, w = alpha.shape
    rgba = np.zeros((h, w, 4), np.uint8); rgba[:, :, :3] = 255; rgba[:, :, 3] = alpha
    Image.fromarray(rgba, "RGBA").save(out)
    os.remove(raw)
    return out


def add_bg(slide, wave=True):
    _ap(slide, _load("bz-light-bg.xml"))
    if wave:
        _ap(slide, _load("bz-wave-layers.xml"))   # 底部 12 层波浪(p98 有 / p99 无)


def ghost_rings(slide, cx, cy, dias):
    for i, d in enumerate(dias):
        al = 13000 - i * 2200
        _ap(slide, parse_xml(
            f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="ghost-ring"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{_e(cx-d/2)}" y="{_e(cy-d/2)}"/><a:ext cx="{_e(d)}" cy="{_e(d)}"/></a:xfrm>'
            f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom><a:noFill/>'
            f'<a:ln w="12700"><a:solidFill><a:srgbClr val="{A1}"><a:alpha val="{al}"/></a:srgbClr></a:solidFill></a:ln></p:spPr></p:sp>'))


def ghost_cluster(slide, x, y, w, h):
    """ghost 圆簇背景(大圆 + 周边小圆,整组一件)—— 调库 bz-ghost-cluster。
    蒸自模板 p99 / 用户在 palette p32 拼的 4 roundRect(并原 ghost-orb+soft-panel+accent-chip 为一件)。
    自然位 (5.14,0.98) ext(7.93×6.04)in。一律调库、禁内联画(§0 步4铁律)。"""
    place(slide, "bz-ghost-cluster.xml", x, y, w, h)


def _text(slide, x, y, w, h, text, sz, bold, color, algn="l", anchor="t", wrap="square"):
    battr = ' b="1"' if bold else ''
    paras = ""
    for line in text.split("\n"):
        paras += (f'<a:p><a:pPr algn="{algn}"><a:lnSpc><a:spcPct val="120000"/></a:lnSpc></a:pPr><a:r><a:rPr lang="zh-CN" sz="{sz}"{battr} dirty="0">'
                  f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                  f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
                  f'<a:t>{escape(line)}</a:t></a:r></a:p>')
    return _ap(slide, parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="t"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="{wrap}" lIns="0" rIns="0" tIns="0" bIns="0" anchor="{anchor}"/>{paras}</p:txBody></p:sp>'))


def visual_slot(slide, x, y, w, h, photo=None, kind="phone", icon="image", adj=14000):
    shadow = ('<a:effectLst><a:outerShdw blurRad="228600" dist="38100" dir="5400000" rotWithShape="0">'
              f'<a:srgbClr val="{A1}"><a:alpha val="24000"/></a:srgbClr></a:outerShdw></a:effectLst>')
    if photo:   # 真图:kind='image'→普通 rect 无白边(车/插图);kind='phone'→roundRect+白边(截图入屏)
        img, rId = slide.part.get_or_add_image_part(photo)
        if kind == "image":
            geom = '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            extra = shadow
        else:
            geom = f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
            extra = f'<a:ln w="25400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>{shadow}'
        return _ap(slide, parse_xml(
            f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="visual-img"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
            f'{geom}{extra}</p:spPr></p:pic>'))
    if kind == "phone":
        return place(slide, "bz-phone-mock.xml", x, y, w, h)   # 调库手机样机(蒸自模板 p99 组合43),禁内联画
    # image 占位(无真图):淡容器脚手架 + 白图标(纯版式占位、非母题;给 photo 即走上面 blipFill 填真图)
    _ap(slide, parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="img-ph"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 8000"/></a:avLst></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{A1}"><a:alpha val="8000"/></a:srgbClr></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="zh-CN"/></a:p></p:txBody></p:sp>'))
    try:
        png = white_icon(icon)
        ic = min(w, h) * 0.30
        img, rId = slide.part.get_or_add_image_part(png)
        _ap(slide, parse_xml(
            f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="img-icon"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{_e(x+w/2-ic/2)}" y="{_e(y+h/2-ic/2)}"/><a:ext cx="{_e(ic)}" cy="{_e(ic)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'))
    except Exception as ex:
        print("  [img-slot icon fallback]", icon, ex)


def feature(slide, cx, cy, label, icon_name, d=0.9):
    _ap(slide, parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="feat-circle"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(cx-d/2)}" y="{_e(cy-d/2)}"/><a:ext cx="{_e(d)}" cy="{_e(d)}"/></a:xfrm>'
        f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'<a:gradFill><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="6FE0FF"/></a:gs><a:gs pos="55000"><a:srgbClr val="{A2}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{A1}"/></a:gs>'
        f'</a:gsLst><a:path path="circle"><a:fillToRect l="30000" t="25000" r="60000" b="70000"/></a:path></a:gradFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:effectLst><a:outerShdw blurRad="76200" dist="19050" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="{A1}"><a:alpha val="32000"/></a:srgbClr></a:outerShdw></a:effectLst></p:spPr></p:sp>'))
    # 真图标(白,从 templates/icons 栅格化),居中叠在圆上
    try:
        png = white_icon(icon_name)
        ic = d * 0.46
        img, rId = slide.part.get_or_add_image_part(png)
        _ap(slide, parse_xml(
            f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="feat-icon"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{_e(cx-ic/2)}" y="{_e(cy-ic/2)}"/><a:ext cx="{_e(ic)}" cy="{_e(ic)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'))
    except Exception as ex:
        print("  [icon fallback]", icon_name, ex)
    _text(slide, cx - 1.05, cy + d / 2 + 0.07, 2.1, 0.5, label, 1400, True, NAVY, algn="ctr", wrap="none")


def build_text_block(slide, label, lead, body, visuals, features=None,
                     rings=(9.4, 3.9, (9.2, 7.0, 4.9, 3.0)), wave=True,
                     clusters=None,
                     lead_y=1.74, body_y=3.30, body_sz=1400):
    add_bg(slide, wave=wave)
    if clusters:                         # ghost 圆簇背景(大圆+周边小圆,整组调库 bz-ghost-cluster)
        for c in clusters:
            ghost_cluster(slide, *c)
    if rings:                            # 同心 ghost 环(p98 用;p99 传 None 跳过)
        ghost_rings(slide, rings[0], rings[1], rings[2])
    for v in visuals:
        visual_slot(slide, *v)
    if features:
        for f in features:
            feature(slide, *f)
    _text(slide, 0.74, 0.46, 5.2, 0.62, label, 2800, True, NAVY, algn="l", wrap="none")
    _text(slide, 0.74, lead_y, 5.6, 1.2, lead, 2800, True, A1, algn="l")
    _text(slide, 0.74, body_y, 4.7, 2.5, body, body_sz, False, BODY, algn="l")


def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_text_block(
        slide,
        label="产品优势",
        lead="一款前沿的智能工具\n专注于跨媒介的创意编排",
        body=("利用最新的人工智能技术，为内容创作者提供一站式的智能创作平台。从灵感构思、"
              "素材整理到多平台分发，全流程智能辅助，让创意表达更高效、更自由。"),
        visuals=[(8.45, 1.75, 1.9, 4.6, None)],
        features=[(9.40, 0.92, "多媒介融合", "share-nodes"),
                  (6.50, 2.00, "实时协作平台", "users"),
                  (12.30, 2.00, "智能内容分析", "chart-line"),
                  (6.50, 5.45, "高级编辑功能", "pen-nib"),
                  (12.30, 5.45, "个性化推荐引擎", "sliders")],
    )
    out = os.path.join(ROOT, "_gallery", "text-block.pptx")
    prs.save(out)
    print("saved", out)


def make_p99_demo(out_path):
    """p99 变体:多视觉簇(车 + 双手机)+ 无特性圈 + 角落 ghost / 柔面板 / 小装饰块 + 无底波浪。"""
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    car = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "rasters", "bz-car-ev.png")
    build_text_block(
        slide,
        label="智能驾驶APP",
        lead="为您提供\n极致的智慧出行体验",
        # 正文 = 模板的 2 个自然段(整段连续、自动换行;**严禁词中间硬断 \n**,上轮被纠)
        body=("采用蓝牙连接车辆，用智能手机控制你的车。相对于传统的显示屏，满足用户十几种模式、一键锁车、中英文等需求。\n"
              "支持手机第三方应用投屏到车机，用哪家导航，车主说了算。与手机原厂合作，体验丝滑的系统级互联体验，越用越爽快。"),
        visuals=[(0.40, 4.55, 5.30, 2.90, car, "image"),    # 车:真高清图 bz-car-ev(plain rect,无白边)
                 (7.55, 2.11, 2.42, 4.91, None, "phone"),    # 手机1:调库 bz-phone-mock
                 (10.14, 0.66, 2.42, 4.91, None, "phone")],  # 手机2:右上对角错位
        features=None, rings=None, wave=False,
        clusters=[(5.138, 0.982, 7.927, 6.037)],             # ghost 圆簇背景(大圆+周边小圆,整组调库 bz-ghost-cluster;自然位)
        lead_y=1.55, body_y=3.04, body_sz=1400,
    )
    prs.save(out_path)
    print("saved", out_path)


if __name__ == "__main__":
    if len(sys.argv) > 2 and sys.argv[1] == "p99":
        make_p99_demo(sys.argv[2])
    else:
        make_demo()
