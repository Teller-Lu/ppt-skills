# -*- coding: utf-8 -*-
"""可复用 团队/人物 网格 builder（business-blue 风格，蒸自 work_inspire_128 p87）。
风格 = 浅蓝渐变底 + airy 白波浪 + 巨大淡水印词 + N 张 zigzag 错位成员卡。
卡骨架 = vectors/bz-team-card.xml（白卡+双幽灵环+头像占位[白色人形]+姓名/职务/分隔/描述，保留 run §6.22）；
背景 = vectors/bz-team-bg.xml（浅蓝渐变底 + airy 白波浪，本页蒸的完整背景）。内容（姓名/职务/描述/照片）用时传入、不写死（§6.24）。
行距：标题单行不套 lnSpc；姓名/职务（单行）与描述（多行 120%）的行距均在 bz-team-card.xml 资产内烤定（design-language §1.3）。

调用：build_team(slide, members, watermark="CORE MEMBER")
  members = [(name, role, desc), ...] 或 (name, role, desc, photo_path)
直接运行本文件 = 用 p87 原内容生成 demo 到 _gallery/people-team.pptx。"""
import os
import sys
import copy
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
VEC = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "vectors")
_parser = etree.XMLParser(resolve_entities=False, no_network=True)


def _load(name):
    return etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)


def _append(slide, el):
    slide.shapes._spTree.append(el)
    return el


# ---------- 背景：浅蓝渐变 + airy 波浪 + 巨大淡水印词 ----------
def add_bg(slide, watermark="CORE MEMBER"):
    bg = _load("bz-team-bg.xml")      # 渐变底 + airy 波浪 合一（本页蒸的完整背景）
    bg.find(f".//{P}cNvPr").set("id", "900")
    _append(slide, bg)
    # 巨大淡水印词（同 §3.2 CONTENTS：sz 166pt 白 alpha~20%）
    word = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="902" name="watermark-word"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="-77724" y="-342900"/><a:ext cx="12636000" cy="2647315"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="none" anchor="t"/><a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="en-US" sz="16600" b="1" dirty="0">'
        f'<a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="20000"/></a:srgbClr></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(watermark)}</a:t></a:r></a:p></p:txBody></p:sp>')
    _append(slide, word)


# ---------- 标题 ----------
def add_title(slide, text):
    tb = slide.shapes.add_textbox(Emu(0), Emu(int(0.45 * IN)), Emu(12192000), Emu(int(0.7 * IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = FACE
    rPr = r._r.get_or_add_rPr()
    for t in ("a:ea", "a:cs"):
        e = rPr.find(qn(t))
        if e is None:
            e = rPr.makeelement(qn(t), {})
            rPr.append(e)
        e.set("typeface", FACE)
    r.font.color.rgb = RGBColor(0x0A, 0x2A, 0x55)
    return tb


def _find(clone, cid):
    for c in clone.iter(f"{P}cNvPr"):
        if c.get("id") == str(cid):
            return c.getparent().getparent()
    return None


def _set_text(clone, cid, text):
    sp = _find(clone, cid)
    if sp is None:
        return
    ts = list(sp.iter(f"{A}t"))
    if ts:
        ts[0].text = text
        for extra in ts[1:]:
            extra.text = ""


def _fill_photo(slide, clone, photo_path):
    # 头像圆(id38) 改填真照片(blipFill) + 去白色人形占位(381/382)
    av = _find(clone, 38)
    if av is None:
        return
    spPr = av.find(f"{P}spPr")
    for tag in ("gradFill", "solidFill", "blipFill", "noFill"):
        e = spPr.find(f"{A}{tag}")
        if e is not None:
            spPr.remove(e)
    img_part, rId = slide.part.get_or_add_image_part(photo_path)
    blip = parse_xml(
        f'<a:blipFill {nsdecls("a", "r")}><a:blip r:embed="{rId}"/>'
        f'<a:stretch><a:fillRect/></a:stretch></a:blipFill>')
    geom = spPr.find(f"{A}prstGeom")
    geom.addnext(blip)   # 椭圆几何把照片裁成圆
    for sid in (381, 382):
        s = _find(clone, sid)
        if s is not None:
            s.getparent().remove(s)


def _renumber(clone, base):
    for i, c in enumerate(clone.iter(f"{P}cNvPr")):
        c.set("id", str(base + i))


# ---------- 主入口 ----------
def build_team(slide, members, watermark="CORE MEMBER",
               card_w=2.708, gap=0.267, y_low=2.367, y_high=1.621):
    """members=[(name,role,desc[,photo]),...]；卡 zigzag 居中铺开，偶序低/奇序高。"""
    N = len(members)
    total = N * card_w + (N - 1) * gap
    x0 = (13.333 - total) / 2.0
    tpl = _load("bz-team-card.xml")
    for k, m in enumerate(members):
        name, role, desc = m[0], m[1], m[2]
        photo = m[3] if len(m) > 3 else None
        clone = copy.deepcopy(tpl)
        _set_text(clone, 26, name)
        _set_text(clone, 27, role)
        _set_text(clone, 28, desc)
        if photo:
            _fill_photo(slide, clone, photo)
        x = x0 + k * (card_w + gap)
        y = y_low if k % 2 == 0 else y_high
        off = clone.find(f"{P}grpSpPr/{A}xfrm/{A}off")
        off.set("x", str(int(x * IN)))
        off.set("y", str(int(y * IN)))
        _renumber(clone, 1000 + k * 100)
        _append(slide, clone)


# ---------- demo：p87 原内容 → _gallery/people-team.pptx ----------
def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "CORE MEMBER")
    add_title(slide, "团队成员介绍")
    members = [
        ("张伟", "技术总监", "负责技术创新和产品开发"),
        ("李丽弦", "运营总监", "管理公司运营和提高效率"),
        ("刘子昂", "财务总监", "监管财务状况和预算规划"),
        ("赵敏", "市场总监", "制定市场策略和品牌推广"),
    ]
    build_team(slide, members, "CORE MEMBER")
    out = os.path.join(ROOT, "_gallery", "people-team.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    make_demo()
