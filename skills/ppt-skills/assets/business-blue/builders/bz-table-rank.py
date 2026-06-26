# -*- coding: utf-8 -*-
"""可复用·参数化表格 builder（business-blue 风格，蒸自 work_inspire_128 p74）。
风格 = 「圆顶角渐变表头底(round2SameRect) + 白卡投影(rect+outerShdw) + 两个透明表格(表头/表体)叠加」。
PPT 表格无法只圆顶两角、也无法整体加圆角+投影 → 故拆成 2 形状 + 2 表格（这是技法精髓）。
全部 srgb（脱主题，免 §6.18 泛色）、每个 run 烤微软雅黑（免 §6.23 回退宋体）、单元格六向边线 noFill（真无框）、无水印。

复用：build_table(slide, x, y, W, headers, rows, ...) 把 4 件加到给定 slide。
内容（数据行/文案）属用时传入，不写死（§6.24 蒸骨架不蒸内容）。
表格单元格 / 标题均单行，按 design-language §1.3 不套 lnSpc（单行无视觉意义；非遗漏）。
直接运行本文件 = 用 p74 原数据生成 demo 到 _gallery/table-rank.pptx。"""
import sys
from xml.sax.saxutils import escape
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
sys.stdout.reconfigure(encoding="utf-8")

IN = 914400  # EMU/inch
FACE = "微软雅黑"
NAVY = "0A2A55"   # tx2 墨字
WHITE = "FFFFFF"  # bg1 表头白字
A1 = "0165FF"     # accent1
A2 = "00BEFF"     # accent2


def _emu(v):
    return Emu(int(round(v * IN)))


def _add_sp(slide, xml):
    el = parse_xml(xml)
    slide.shapes._spTree.append(el)
    return el


# ---------- 背景（满版浅蓝竖渐变，demo 用，非表格件） ----------
def add_bg(slide, _id=900):
    xml = (
        f'<p:sp {nsdecls("p", "a")}>'
        f'<p:nvSpPr><p:cNvPr id="{_id}" name="bg"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A1}"><a:lumMod val="2000"/><a:lumOff val="98000"/></a:srgbClr></a:gs>'
        f'<a:gs pos="88000"><a:srgbClr val="{A1}"><a:lumMod val="8000"/><a:lumOff val="92000"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="5400000" scaled="0"/><a:tileRect/></a:gradFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    return _add_sp(slide, xml)


# ---------- 标题（demo 用，非表格件） ----------
def add_title(slide, text, _id=903):
    tb = slide.shapes.add_textbox(_emu(0.742), _emu(0.46), _emu(10.4), _emu(0.7))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.name = FACE
    rPr = r._r.get_or_add_rPr()
    for t in ("a:ea", "a:cs"):
        e = rPr.find(qn(t))
        if e is None:
            e = rPr.makeelement(qn(t), {})
            rPr.append(e)
        e.set("typeface", FACE)
    r.font.color.rgb = RGBColor(0x0A, 0x2A, 0x55)
    return tb


# ---------- 表头底：圆顶角 + 径向渐变（accent2@右下中心 → accent1@左上边） ----------
def add_header_bar(slide, x, y, cx, cy, _id=901):
    xml = (
        f'<p:sp {nsdecls("p", "a")}>'
        f'<p:nvSpPr><p:cNvPr id="{_id}" name="tbl-header-bar"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(x*IN)}" y="{int(y*IN)}"/><a:ext cx="{int(cx*IN)}" cy="{int(cy*IN)}"/></a:xfrm>'
        f'<a:prstGeom prst="round2SameRect"><a:avLst/></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="15000"><a:srgbClr val="{A2}"/></a:gs>'
        f'<a:gs pos="87000"><a:srgbClr val="{A1}"/></a:gs>'
        f'</a:gsLst><a:path path="circle"><a:fillToRect r="100000" b="100000"/></a:path>'
        f'<a:tileRect l="-100000" t="-100000"/></a:gradFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    return _add_sp(slide, xml)


# ---------- 白卡：纯白 rect + 向下柔蓝投影 ----------
def add_card(slide, x, y, cx, cy, _id=902):
    xml = (
        f'<p:sp {nsdecls("p", "a")}>'
        f'<p:nvSpPr><p:cNvPr id="{_id}" name="tbl-card"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(x*IN)}" y="{int(y*IN)}"/><a:ext cx="{int(cx*IN)}" cy="{int(cy*IN)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{WHITE}"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:effectLst><a:outerShdw blurRad="469900" dist="342900" dir="5400000" sx="96000" sy="96000" algn="t" rotWithShape="0">'
        f'<a:srgbClr val="{A1}"><a:lumMod val="75000"/><a:alpha val="40000"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst></p:spPr></p:sp>')
    return _add_sp(slide, xml)


# ---------- 表格样式：去默认 firstRow/band，套「无样式无网格」 ----------
def _style_table(gf):
    tbl = gf._element.find(".//" + qn("a:tbl"))
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = parse_xml(f'<a:tblPr {nsdecls("a")}/>')
        tbl.insert(0, tblPr)
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is None:
        sid = parse_xml(f'<a:tableStyleId {nsdecls("a")}>x</a:tableStyleId>')
        tblPr.append(sid)
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid


# ---------- 单元格：透明六向边线 + 填充 + 一段一 run（烤雅黑） ----------
def _set_cell(tc, text, sz, bold, color, marL, marR=None, band=False, band_alpha=50000):
    # txBody：清旧段，写一段一 run，烤 latin/ea/cs
    txBody = tc.find(qn("a:txBody"))
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    battr = ' b="1"' if bold else ""
    p = parse_xml(
        f'<a:p {nsdecls("a")}><a:r>'
        f'<a:rPr lang="zh-CN" altLang="en-US" sz="{sz}"{battr} dirty="0">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/>'
        f'</a:rPr><a:t>{escape(text)}</a:t></a:r></a:p>')
    txBody.append(p)
    # tcPr：六向边线 noFill（真无框）+ 填充
    old = tc.find(qn("a:tcPr"))
    if old is not None:
        tc.remove(old)
    if band:
        fillxml = (f'<a:solidFill><a:srgbClr val="{A1}">'
                   f'<a:lumMod val="20000"/><a:lumOff val="80000"/><a:alpha val="{band_alpha}"/>'
                   f'</a:srgbClr></a:solidFill>')
    else:
        fillxml = "<a:noFill/>"
    marRattr = f' marR="{marR}"' if marR is not None else ""
    tcPr = parse_xml(
        f'<a:tcPr {nsdecls("a")} marL="{marL}"{marRattr} anchor="ctr">'
        f'<a:lnL w="12700" cmpd="sng"><a:noFill/></a:lnL>'
        f'<a:lnR w="12700" cmpd="sng"><a:noFill/></a:lnR>'
        f'<a:lnT w="12700" cmpd="sng"><a:noFill/></a:lnT>'
        f'<a:lnB w="12700" cmpd="sng"><a:noFill/></a:lnB>'
        f'<a:lnTlToBr w="12700" cmpd="sng"><a:noFill/></a:lnTlToBr>'
        f'<a:lnBlToTr w="12700" cmpd="sng"><a:noFill/></a:lnBlToTr>'
        f'{fillxml}</a:tcPr>')
    tc.append(tcPr)


def _col_widths(N, Wt, col_ratios):
    if col_ratios is None:
        col_ratios = [802372, 1839675] + [1917588] * (N - 2)
    s = sum(col_ratios)
    widths = [int(Wt * IN * r / s) for r in col_ratios]
    widths[-1] = int(Wt * IN) - sum(widths[:-1])  # 余量归末列，免累积误差
    return widths


# ---------- 主入口：把表头底 + 白卡 + 表头表 + 表体表 加到 slide ----------
def build_table(slide, x, y, W, headers, rows, col_ratios=None,
                bold_cols=(2,), band_alpha=50000, header_sz=1600, body_sz=1400,
                inset=0.189, header_h=0.611, row_h=0.525):
    """x/y = 表头底左上角(in)，W = 表头底/白卡宽(in)。表格在其内缩进 inset。
    bold_cols = 要整列加粗的列号（强调主指标列，模板=第3列里程）。"""
    N = len(headers)
    M = len(rows)
    xt = x + inset
    Wt = W - 2 * inset
    widths = _col_widths(N, Wt, col_ratios)
    header_tbl_y = y + (header_h - 0.517) / 2.0     # 表头表在表头底内垂直居中
    body_tbl_y = y + header_h + 0.04                # 表体表紧贴表头底下方
    body_tbl_h = M * row_h
    card_y = body_tbl_y - 0.113                      # 白卡上沿塞到表头底之下
    card_h = body_tbl_h + 0.15                        # 上 0.113 + 下 0.037 包住表体

    # z 序：白卡 → 表体表 → 表头底 → 表头表
    add_card(slide, x, card_y, W, card_h)
    # 表体表
    gf = slide.shapes.add_table(M, N, _emu(xt), _emu(body_tbl_y), _emu(Wt), _emu(body_tbl_h))
    _style_table(gf)
    tb = gf.table
    for i, w in enumerate(widths):
        tb.columns[i].width = Emu(w)
    for rr in tb.rows:
        rr.height = _emu(row_h)
    for ri, row in enumerate(rows):
        band = (ri % 2 == 1)
        for ci, val in enumerate(row):
            marL = 288000 if ci == 0 else 252000
            marR = 72000 if ci == 0 else None
            _set_cell(tb.cell(ri, ci)._tc, str(val), body_sz, ci in bold_cols,
                      NAVY, marL, marR, band=band, band_alpha=band_alpha)
    # 表头底
    add_header_bar(slide, x, y, W, header_h)
    # 表头表
    gfh = slide.shapes.add_table(1, N, _emu(xt), _emu(header_tbl_y), _emu(Wt), _emu(0.517))
    _style_table(gfh)
    th = gfh.table
    for i, w in enumerate(widths):
        th.columns[i].width = Emu(w)
    th.rows[0].height = _emu(0.517)
    for ci, h in enumerate(headers):
        marL = 180000 if ci == 0 else 252000
        _set_cell(th.cell(0, ci)._tc, h, header_sz, True, WHITE, marL, None, band=False)
    return gf, gfh


# ---------- demo：p74 原数据，生成 _gallery/table-rank.pptx ----------
def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "自动驾驶测试数据排名")
    headers = ["排名", "公司", "测试总里程(英里)", "脱离次数", "车辆数量(台)", "MPD"]
    rows = [
        ["1", "小马智行", "305616.73", "21", "38", "14553.18"],
        ["2", "文远知行", "57966.25", "3", "14", "19332.08"],
        ["3", "AutoX", "50108", "1", "44", "56108"],
        ["4", "滴滴", "40744.67", "1", "12", "40744.67"],
        ["5", "元戎启行", "30872", "2", "2", "15436"],
        ["6", "百度Apollo", "9564.6", "0", "3", "无脱离"],
        ["7", "小马智行", "9042.87", "1", "3", "9042.87"],
        ["8", "轻舟智航", "6320", "5", "2", "1264.00"],
        ["9", "百度Apollo", "1467.5", "1", "5", "1467.50"],
    ]
    build_table(slide, x=0.839, y=1.574, W=11.656, headers=headers, rows=rows,
                bold_cols=[2], band_alpha=50000)
    out = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills\_gallery\table-rank.pptx"
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    make_demo()
