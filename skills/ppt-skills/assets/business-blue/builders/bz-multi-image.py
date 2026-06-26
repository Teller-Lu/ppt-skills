# -*- coding: utf-8 -*-
"""可复用 多图排版 multi-image builder（business-blue，蒸自 work_inspire_128 p103-106）。
本批首个版式 = p106「左文 + 右手机簇」：页标题 28 + slogan 36 + 2 图标特性条 + 右 N 手机样机错落 + 底部蓝带。
手机调库 bz-phone-mock、底调库 bz-light-bg；图标特性条 / 手机簇布局属纯版式、builder 现写（§资产架构:母题入库、纯版式可现写）。
图标取 templates/icons 真图标栅格化成白透；截图 / 文字属内容、用时填（§6.24）。

行距：标题 / slogan / 特性名均单行，特性说明单行——按 design-language §1.3 单行不强求 lnSpc（_text 通用，多行才有意义）。

调用：build_p106(slide, title, slogan, features, phones)
  features=[(name, desc, icon_name), ...]        phones=[(x,y,w,h,screenshot_or_None), ...]
直接运行 = p106 内容 demo → _experiments/p106_design/p106.pptx。"""
import os
import sys
import tempfile
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
A1 = "0165FF"
A2 = "00BEFF"
NAVY = "0A2A55"
BODY = "44546A"
WHITE = "FFFFFF"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
VEC = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "vectors")
RASTERS = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "rasters")
ICONLIB = r"D:\Lu.Yao7\github\ppt-master\skills\ppt-master\templates\icons"
ICACHE = os.path.join(tempfile.gettempdir(), "bz_iconcache")
os.makedirs(ICACHE, exist_ok=True)
_parser = etree.XMLParser(resolve_entities=False, no_network=True)
_uid = [7000]


def _nid():
    _uid[0] += 1
    return _uid[0]


def _e(v):
    return int(round(v * IN))


def _ap(slide, el):
    slide.shapes._spTree.append(el)
    return el


def _load(name):
    el = etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    return el


def place(slide, name, x, y, w, h):
    """从素材库载入资产并定位/缩放——禁内联画图,一律调库(§资产架构)。renumber 全 cNvPr 避免多实例子件 id 相撞。"""
    el = etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)
    for nv in el.iter(f"{P}cNvPr"):
        nv.set("id", str(_nid()))
    xf = el.find(f"{P}grpSpPr/{A}xfrm") if el.tag == f"{P}grpSp" else el.find(f"{P}spPr/{A}xfrm")
    xf.find(f"{A}off").set("x", str(_e(x))); xf.find(f"{A}off").set("y", str(_e(y)))
    xf.find(f"{A}ext").set("cx", str(_e(w))); xf.find(f"{A}ext").set("cy", str(_e(h)))
    return _ap(slide, el)


def white_icon(name, sub="chunk-filled", scale=8):
    """templates/icons/<sub>/<name>.svg → 白色透明 PNG(缓存),供 add_picture 叠到图标块。"""
    out = os.path.join(ICACHE, f"{sub}_{name}_w.png")
    if os.path.exists(out):
        return out
    from svglib.svglib import svg2rlg
    from reportlab.graphics import renderPM
    import numpy as np
    from PIL import Image
    d = svg2rlg(os.path.join(ICONLIB, sub, name + ".svg"))
    d.scale(scale, scale); d.width *= scale; d.height *= scale
    raw = out + "_raw.png"
    renderPM.drawToFile(d, raw, fmt="PNG", bg=0xFFFFFF)
    g = np.array(Image.open(raw).convert("L")).astype(np.int16)
    alpha = np.clip(255 - g, 0, 255).astype(np.uint8)
    h, w = alpha.shape
    rgba = np.zeros((h, w, 4), np.uint8); rgba[:, :, :3] = 255; rgba[:, :, 3] = alpha
    Image.fromarray(rgba, "RGBA").save(out)
    os.remove(raw)
    return out


def _text(slide, x, y, w, h, text, sz, bold, color, algn="l", anchor="t", wrap="square"):
    """通用文字:多行按 \n 分段、统一 lnSpc 120%(单行加亦无害);烤微软雅黑 latin/ea/cs(§6.23)。"""
    battr = ' b="1"' if bold else ''
    paras = ""
    for line in text.split("\n"):
        paras += (f'<a:p><a:pPr algn="{algn}"><a:lnSpc><a:spcPct val="120000"/></a:lnSpc></a:pPr><a:r><a:rPr lang="zh-CN" sz="{sz}"{battr} dirty="0">'
                  f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                  f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
                  f'<a:t>{escape(line)}</a:t></a:r></a:p>')
    return _ap(slide, parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="t"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="{wrap}" lIns="0" rIns="0" tIns="0" bIns="0" anchor="{anchor}"/>{paras}</p:txBody></p:sp>'))


def visual_slot(slide, x, y, w, h, photo=None, kind="phone", adj=14000):
    """手机/图片视觉槽:kind='phone'→调库 bz-phone-mock(空屏占位);给 photo→roundRect blipFill 入截图。"""
    shadow = ('<a:effectLst><a:outerShdw blurRad="228600" dist="38100" dir="5400000" rotWithShape="0">'
              f'<a:srgbClr val="{A1}"><a:alpha val="24000"/></a:srgbClr></a:outerShdw></a:effectLst>')
    if photo:
        img, rId = slide.part.get_or_add_image_part(photo)
        geom = f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
        extra = f'<a:ln w="25400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>{shadow}'
        return _ap(slide, parse_xml(
            f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="phone-shot"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
            f'{geom}{extra}</p:spPr></p:pic>'))
    return place(slide, "bz-phone-mock.xml", x, y, w, h)   # 空屏手机:调库,禁内联画


def add_bg(slide):
    _ap(slide, _load("bz-light-bg.xml"))


def bottom_band(slide):
    """底部 = 蓝带(调库 bz-bottom-blue) + 白波浪(调库 bz-bottom-wave,盖蓝带成波浪上沿)。原模板双层,零内联。"""
    place(slide, "bz-bottom-blue.xml", 0, 6.34, 13.33, 1.16)
    place(slide, "bz-bottom-wave.xml", 0, 6.07, 13.33, 1.43)


def icon_feature_bar(slide, ix, iy, name, desc, icon_vec, box=0.64):
    """图标特性条 = 调库白描边图标框 + 调库原生图标 + 特性名(18粗) + 说明(16灰)。零内联。"""
    place(slide, "bz-icon-box.xml", ix, iy, box, box)                  # 白描边框(调库)
    place(slide, icon_vec, ix + box * 0.18, iy + box * 0.18, box * 0.64, box * 0.64)  # 原生图标(调库)
    tx = ix + box + 0.18
    _text(slide, tx, iy - 0.17, 3.4, 0.42, name, 1800, True, NAVY, wrap="none")
    _text(slide, tx, iy + 0.22, 3.6, 0.42, desc, 1600, False, BODY, wrap="none")


def build_p106(slide, title, slogan, features, phones):
    place(slide, "bz-social-bg.xml", 0, 0, 13.33, 7.5)   # p106 背景(底部偏蓝,调库)
    bottom_band(slide)
    for ph in phones:                       # 手机簇(列表顺序=z序,后者压前者)
        visual_slot(slide, ph[0], ph[1], ph[2], ph[3], ph[4], kind="phone")
    _text(slide, 1.04, 1.20, 5.0, 0.60, title, 2800, True, NAVY, wrap="none")     # 页标题 28
    _text(slide, 1.04, 1.84, 4.9, 1.40, slogan, 3600, True, NAVY)                  # slogan 36(可断行)
    for i, (name, desc, icon) in enumerate(features):                             # 2 图标特性条
        icon_feature_bar(slide, 1.12, 4.02 + i * 1.20, name, desc, icon)


# flow_podium 已废弃 → build_p103 改调库 bz-flow-podium-trio.xml（原模板真件:3联倒梯形+scene3d透视+组渐变）


# tag_bar 已废弃 → build_p103 改调库 bz-tag-bar.xml（原模板真件:径向渐变胶囊+倒影）


def build_p103(slide, title, columns):
    """三屏横排:标题 + 三联承托台(调库) + 3 列[大屏图 + 向下箭头(调库) + 小标题20 + 正文15 + 标签条(调库)]。
    columns=[(image_or_None, heading, body, tag), ...] 取前 3 列。母题全调库、零内联(蒸馏流程§步4铁律)。"""
    add_bg(slide)
    place(slide, "bz-flow-podium-trio.xml", -0.18, 2.34, 13.69, 2.47)   # 三联承托台(整组:倒梯形+scene3d透视+组渐变)
    xs = [1.09, 5.20, 9.18]
    arrow_x = [2.23, 6.52, 10.87]                                       # 箭头列间(原模板位)
    IMG_W, IMG_H = 2.96, 1.67
    for k, (img, heading, body, tag) in enumerate(columns[:3]):
        cx = xs[k]
        visual_slot(slide, cx, 2.02, IMG_W, IMG_H, img, adj=5000)       # 大屏图(白边圆角截图卡)
        place(slide, "bz-down-arrow.xml", arrow_x[k], 5.62, 0.29, 0.30)  # 向下箭头(调库)
        _text(slide, cx + 0.05, 4.29, 2.80, 0.44, heading, 2000, True, A1, wrap="none")   # 小标题20
        _text(slide, cx + 0.02, 4.84, 2.82, 1.00, body, 1500, False, BODY)                # 正文15
        place(slide, "bz-tag-bar.xml", cx - 0.05, 6.11, 2.93, 0.53)     # 胶囊标签条(调库)
        _text(slide, cx - 0.05, 6.17, 2.93, 0.42, tag, 1400, False, WHITE, algn="ctr", anchor="ctr", wrap="none")  # 标签文字
    _text(slide, 0.74, 0.49, 10.4, 0.63, title, 2800, True, NAVY, wrap="none")            # 页标题28


def make_p103_demo(out_path=None):
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sd = os.path.join(ROOT, "_experiments", "p103_design")
    build_p103(
        slide,
        title="基于“一张图”的生态环境管理",
        columns=[
            (os.path.join(sd, "screen1.png"), "“一张图”集成展示",
             "缺乏基于“一张图”的跨区域、跨业务可视化总览", "实现生态环境信息“一览无余”"),
            (os.path.join(sd, "screen2.png"), "“一张图”决策管理",
             "缺乏信息内在关系、发展规律的深入挖掘和决策分析", "实现跨业务综合研判辅助决策"),
            (os.path.join(sd, "screen3.png"), "“一张图”共享服务",
             "缺乏统一的空间信息技术支撑服务，存在重复建设", "资金集约使用，提升服务质量"),
        ],
    )
    out = out_path or os.path.join(sd, "p103.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("saved", out)


def add_p105_bg(slide):
    """p105 背景:与 p104 同一隧道透视图(duotone染蓝+29%透明,还原原模板[39]=[22]同 blob,否则发白)。"""
    pic = slide.shapes.add_picture(os.path.join(RASTERS, "bz-tilt-bg.jpg"), 0, 0, _e(13.33), _e(7.5))
    blip = pic._element.find(f".//{A}blip")
    blip.append(parse_xml(f'<a:duotone {nsdecls("a")}><a:srgbClr val="0165FF"><a:shade val="45000"/><a:satMod val="135000"/></a:srgbClr><a:prstClr val="white"/></a:duotone>'))
    blip.append(parse_xml(f'<a:alphaModFix {nsdecls("a")} amt="29000"/>'))


def corner_tile(slide, x, y, w, h, img, label):
    """四角场景小图(白边圆角) + 调库半透明蓝标签条(底部) + 白字16粗。"""
    visual_slot(slide, x, y, w, h, img, adj=6000)
    bh = 0.6
    place(slide, "bz-corner-tag.xml", x, y + h - bh, w, bh)   # 标签条(调库)
    _text(slide, x, y + h - bh + 0.12, w, bh - 0.2, label, 1600, True, WHITE, algn="ctr", anchor="ctr", wrap="none")


def build_p105(slide, title, center_img, system_name, corners):
    """中心辐射:隧道透视背景 + 中心大图(坐透视底座) + 两摄像头 + FPG系统名 + 四角卫星小图(带标签)。
    corners=[(img, label), ...] 顺序 左上/右上/左下/右下。母题全调库、零内联(蒸馏流程§步4铁律)。"""
    add_p105_bg(slide)                                                  # 隧道透视背景(与 p104 同图,duotone染蓝半透明)
    place(slide, "bz-podium-3d.xml", 3.49, 4.12, 6.36, 3.34)           # 中心透视底座(梯形+反光+scene3d,调库)
    visual_slot(slide, 4.52, 2.62, 4.29, 2.83, center_img, adj=6000)    # 中心大图
    slide.shapes.add_picture(os.path.join(RASTERS, "bz-camera.png"), _e(5.78), _e(4.59), _e(0.93), _e(1.17))  # 摄像头左
    slide.shapes.add_picture(os.path.join(RASTERS, "bz-camera.png"), _e(6.68), _e(4.59), _e(0.93), _e(1.17))  # 摄像头右
    _text(slide, 4.66, 6.40, 4.6, 0.5, system_name, 2400, True, A1, algn="ctr", wrap="none")  # FPG系统名24
    pos = [(1.35, 1.69, 2.69, 1.78), (9.29, 1.46, 2.85, 1.88),
           (0.63, 4.08, 2.69, 1.78), (10.01, 3.96, 3.0, 1.98)]          # 四角
    for (px, py, pw, ph), (img, label) in zip(pos, corners):
        corner_tile(slide, px, py, pw, ph, img, label)
    _text(slide, 0.74, 0.49, 10.4, 0.63, title, 2800, True, NAVY, wrap="none")   # 页标题28


def make_p105_demo(out_path=None):
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sd = os.path.join(ROOT, "_experiments", "p105_design")
    build_p105(
        slide,
        title="生产信息集中管理，优化制程并降低成本",
        center_img=os.path.join(sd, "center.png"),
        system_name="FPG 智能监控系统",
        corners=[(os.path.join(sd, "corner1.png"), "生产管理"),
                 (os.path.join(sd, "corner2.png"), "设备管理"),
                 (os.path.join(sd, "corner3.png"), "预测管理"),
                 (os.path.join(sd, "corner4.png"), "能耗管理")],
    )
    out = out_path or os.path.join(sd, "p105.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("saved", out)


def add_dark_bg(slide):
    _ap(slide, _load("bz-dark-bg.xml"))


def add_tilt_bg(slide):
    """p104 背景:浅蓝渐变底([11]调库) + 隧道图(duotone染蓝+33%透明,还原原模板[22]) + 透视光带 + 角光束。"""
    place(slide, "bz-tilt-base.xml", 0, 0, 13.33, 7.5)              # 浅蓝渐变底
    pic = slide.shapes.add_picture(os.path.join(RASTERS, "bz-tilt-bg.jpg"), 0, 0, _e(13.33), _e(7.5))
    blip = pic._element.find(f".//{A}blip")                          # 还原原模板 duotone 染蓝 + 半透明(否则发白)
    blip.append(parse_xml(f'<a:duotone {nsdecls("a")}><a:srgbClr val="0165FF"><a:shade val="45000"/><a:satMod val="135000"/></a:srgbClr><a:prstClr val="white"/></a:duotone>'))
    blip.append(parse_xml(f'<a:alphaModFix {nsdecls("a")} amt="33000"/>'))
    place(slide, "bz-tilt-bg-glow.xml", -2.5, -0.88, 18.33, 10.0)    # 背景透视光带
    place(slide, "bz-tilt-corner-glow.xml", 1.69, 4.45, 2.7, 1.71)   # 左下角光束
    place(slide, "bz-tilt-corner-glow.xml", 8.9, 4.92, 2.7, 1.71)    # 右下角光束


def tilt_card(slide, x, y, w, h, img, label, label_sz=2000):
    """倾斜卡:snip2DiagRect(剪角矩形)裁图+白边+倒影 + 调库蓝发光层 + 调库渐变边框 + 卡内底白标签。
    形状/发光层/边框均蒸自原模板 p104,发光层+边框调库(bz-tilt-glow/frame),零内联 substantial 母题。"""
    img_part, rId = slide.part.get_or_add_image_part(img)
    _ap(slide, parse_xml(
        f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="tilt-img"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="snip2DiagRect"><a:avLst><a:gd name="adj1" fmla="val 11586"/><a:gd name="adj2" fmla="val 10606"/></a:avLst></a:prstGeom>'
        f'<a:ln w="19050"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>'
        f'<a:effectLst><a:reflection blurRad="177800" stA="52000" endA="300" endPos="23000" dir="5400000" sy="-100000" algn="bl" rotWithShape="0"/></a:effectLst></p:spPr></p:pic>'))
    place(slide, "bz-tilt-glow.xml", x, y, w, h)       # 蓝发光层(调库)
    place(slide, "bz-tilt-frame.xml", x, y, w, h)      # 渐变边框(调库)
    _text(slide, x, y + h - 0.52, w, 0.46, label, label_sz, True, WHITE, algn="ctr", anchor="ctr", wrap="none")


def build_p104(slide, title, cards):
    """4 张倾斜卡(剪角矩形裁图+蓝发光层+渐变边框,全调库)错落 + 暗蓝光带背景 + 白标题。
    cards=[(x,y,w,h,img,label,sz), ...]。母题全调库、零内联(蒸馏流程§步4铁律)。"""
    add_tilt_bg(slide)
    for c in cards:
        tilt_card(slide, *c)
    _text(slide, 0.74, 0.49, 10.4, 0.63, title, 2800, True, WHITE, wrap="none")


def make_p104_demo(out_path=None):
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sd = os.path.join(ROOT, "_experiments", "p104_design")
    build_p104(
        slide,
        title="应急调度",
        # 错落:分析(左)/系统联动(中上)/可视化(右)/一键式调度(中下大,最后画压前)
        cards=[(0.83, 2.47, 3.54, 2.02, os.path.join(sd, "card1.png"), "分析研判", 2000),
               (5.20, 1.34, 3.23, 1.84, os.path.join(sd, "card2.png"), "系统联动", 2000),
               (8.97, 2.47, 3.54, 2.02, os.path.join(sd, "card4.png"), "可视化处置", 2000),
               (4.66, 3.69, 4.01, 2.29, os.path.join(sd, "card3.png"), "一键式调度", 2400)],
    )
    out = out_path or os.path.join(sd, "p104.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("saved", out)


def make_p106_demo(out_path=None):
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_p106(
        slide,
        title="自媒体账号一键管理",
        slogan="品宣、矩阵、引流\n必备神器",
        features=[("一键管理", "近 50+ 主流自媒体账号", "bz-icon-users.xml"),
                  ("一键发布", "至 30+ 主流平台", "bz-icon-share.xml")],
        # 手机簇:左/右先画、中间最高最后画压在前(参考非照抄,自排错落);填 demo 示意截图(用时换真截图)
        phones=[(6.55, 2.30, 2.25, 4.70, os.path.join(ROOT, "_experiments", "p106_design", "shot1.png")),
                (10.55, 2.05, 2.25, 4.70, os.path.join(ROOT, "_experiments", "p106_design", "shot3.png")),
                (8.45, 1.15, 2.55, 5.30, os.path.join(ROOT, "_experiments", "p106_design", "shot2.png"))],
    )
    out = out_path or os.path.join(ROOT, "_experiments", "p106_design", "p106.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    mode = sys.argv[1] if len(sys.argv) > 1 else "p106"
    {"p103": make_p103_demo, "p104": make_p104_demo, "p105": make_p105_demo, "p106": make_p106_demo}.get(mode, make_p106_demo)()
