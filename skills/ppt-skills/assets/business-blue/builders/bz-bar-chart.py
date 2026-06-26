# -*- coding: utf-8 -*-
"""可复用 数据图表(柱状)builder（business-blue，蒸自 work_inspire_128 p117）。
柱图**用形状重绘**（渐变柱 + 数值标签 + 类别标签 + 趋势箭头），非原生 chart——本 skill 原生可控哲学、
免内嵌 Excel、完全可控。缺口顶白卡 = vectors/bz-stat-card.xml（本页蒸的）。数据用时传（§6.24）。

调用：
  bar_chart(slide, x, y, w, h, cats, vals, hl_idx=-1, hl="a1", arrow=None)   # 单个柱图
  stat_card(slide, cx, date, kpis, charts)                                    # 缺口卡 + KPI + N图
所有 _textbox 为 wrap=none 单行精调标签（位置像素级控），按 design-language §1.3 严禁套 lnSpc（anchor=t 单行套 120% 会基线下移、破坏精调；非遗漏）。
直接运行 = p117 三卡 demo → _gallery/data-chart.pptx。"""
import os
import sys
import copy
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
NAVY = "0A2A55"
GRAY = "98A2B3"
A1 = "0165FF"
A2 = "00BEFF"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
VEC = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "vectors")
_parser = etree.XMLParser(resolve_entities=False, no_network=True)
_uid = [2000]

# 渐变柱配色：top→bottom
GRADS = {
    "gray": ("E4E8EE", "BFC7D3"),
    "a1": ("4D8DFF", "0165FF"),
    "a2": ("5BD4FF", "00BEFF"),
}
HLLABEL = {"a1": "0165FF", "a2": "00AEE6", "gray": GRAY}


def _nid():
    _uid[0] += 1
    return _uid[0]


def _load(name):
    return etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)


def _append(slide, el):
    slide.shapes._spTree.append(el)
    return el


def _e(v):
    return int(round(v * IN))


def _textbox(slide, x, y, w, h, text, sz, bold, color, algn="ctr", anchor="t"):
    battr = ' b="1"' if bold else ''
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="t"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="none" lIns="0" rIns="0" tIns="0" bIns="0" anchor="{anchor}"/>'
        f'<a:p><a:pPr algn="{algn}"/><a:r><a:rPr lang="zh-CN" sz="{sz}"{battr} dirty="0">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>')
    return _append(slide, sp)


def _bar(slide, x, y, w, h, grad):
    top, bot = GRADS[grad]
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="bar"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="round2SameRect"><a:avLst><a:gd name="adj1" fmla="val 14000"/><a:gd name="adj2" fmla="val 0"/></a:avLst></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{top}"/></a:gs><a:gs pos="100000"><a:srgbClr val="{bot}"/></a:gs>'
        f'</a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill><a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    return _append(slide, sp)


def trend_arrow(slide, x0, y0, x1, y1, grad, down=False):
    ax, ay = min(x0, x1), min(y0, y1)
    W, H = max(abs(x1 - x0), 0.05), max(abs(y1 - y0), 0.05)
    we, he = _e(W), _e(H)
    if not down:   # 上行：左下→右上
        path = f'<a:moveTo><a:pt x="0" y="{he}"/></a:moveTo><a:quadBezTo><a:pt x="{int(we*0.5)}" y="{int(he*0.62)}"/><a:pt x="{we}" y="0"/></a:quadBezTo>'
    else:          # 下行：左上→右下
        path = f'<a:moveTo><a:pt x="0" y="0"/></a:moveTo><a:quadBezTo><a:pt x="{int(we*0.5)}" y="{int(he*0.4)}"/><a:pt x="{we}" y="{he}"/></a:quadBezTo>'
    col = A1 if grad == "a1" else (A2 if grad == "a2" else GRAY)
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="trend"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(ax)}" y="{_e(ay)}"/><a:ext cx="{we}" cy="{he}"/></a:xfrm>'
        f'<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/><a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst><a:path w="{we}" h="{he}">{path}</a:path></a:pathLst></a:custGeom>'
        f'<a:noFill/><a:ln w="28575" cap="rnd"><a:solidFill><a:srgbClr val="{col}"><a:alpha val="90000"/></a:srgbClr></a:solidFill>'
        f'<a:tailEnd type="triangle" w="med" len="lg"/></a:ln></p:spPr></p:sp>')
    return _append(slide, sp)


def bar_chart(slide, x, y, w, h, cats, vals, hl_idx=-1, hl="a1", arrow=None, note=None):
    """在 (x,y,w,h) 内画柱图：渐变柱 + 数值(上) + 类别(下) + 可选趋势箭头 + 可选标注。
    hl_idx=高亮柱序(-1 末柱,None 无高亮)；hl∈{a1,a2}；arrow∈{None,'up','down'}；note=箭头旁标注文字。"""
    n = len(vals)
    top_pad, bot_pad = 0.34, 0.30
    plot_h = h - top_pad - bot_pad
    base_y = y + top_pad + plot_h
    slot = w / n
    bw = slot * 0.5
    mx = max(vals) or 1
    if hl_idx is not None and hl_idx < 0:
        hl_idx = n + hl_idx
    tops = []
    for i, v in enumerate(vals):
        bh = (v / mx) * plot_h
        bx = x + i * slot + (slot - bw) / 2
        bt = base_y - bh
        tops.append((bx + bw / 2, bt))
        grad = hl if (hl_idx is not None and i == hl_idx) else "gray"
        _bar(slide, bx, bt, bw, bh, grad)
        # 数值标签
        is_hl = (hl_idx is not None and i == hl_idx)
        vs = ("%g" % v)
        _textbox(slide, x + i * slot, bt - 0.30, slot, 0.28, vs,
                 1400 if is_hl else 1200, is_hl, HLLABEL[hl] if is_hl else GRAY)
        # 类别标签
        _textbox(slide, x + i * slot, base_y + 0.04, slot, 0.26, str(cats[i]), 1100, False, GRAY)
    # 基线
    base = parse_xml(
        f'<p:cxnSp {nsdecls("p", "a")}><p:nvCxnSpPr><p:cNvPr id="{_nid()}" name="base"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(base_y)}"/><a:ext cx="{_e(w)}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="D5DAE2"/></a:solidFill></a:ln></p:spPr></p:cxnSp>')
    _append(slide, base)
    # 趋势箭头（优雅曲线，走柱体/空白、不压数值标签）
    if arrow == "up":
        trend_arrow(slide, x + slot * 0.42, base_y - plot_h * 0.16,
                    tops[-1][0], tops[-1][1] + 0.05, hl, down=False)
    elif arrow == "down":
        trend_arrow(slide, tops[0][0], tops[0][1] + 0.40,
                    tops[-1][0], tops[-1][1] + 0.74, "gray", down=True)
    # 标注（如「下降4.78%」）：压柱体上方、数值标签之下，1200 粗，不遮数字
    if note:
        _textbox(slide, x + w * 0.16, y + top_pad + 0.14, w * 0.68, 0.5, note, 1200, True, NAVY, algn="ctr")


# ---------- 缺口卡 + KPI + N 图 ----------
def stat_card(slide, cx, date, kpis, charts, card_y=1.89):
    card = _load("bz-stat-card.xml")
    card.find(f".//{P}cNvPr").set("id", str(_nid()))
    off = card.find(f"{P}grpSpPr/{A}xfrm/{A}off")
    off.set("x", str(_e(cx))); off.set("y", str(_e(card_y)))
    _append(slide, card)
    pad = 0.12
    _textbox(slide, cx + pad, card_y + 0.36, 2.0, 0.28, date, 1100, False, GRAY, algn="l")
    # 两列 KPI：label 上 + value 下
    kx = [cx + pad, cx + pad + 2.02]
    for j, (lab, val, vcol) in enumerate(kpis):
        _textbox(slide, kx[j], card_y + 0.69, 1.9, 0.32, lab, 1400, False, NAVY, algn="l")
        _textbox(slide, kx[j], card_y + 0.98, 1.9, 0.42, val, 1800, True, vcol, algn="l")
    # 图（每个含类别题 + 柱图）
    for spec in charts:
        _textbox(slide, spec["x"] + 0.02, card_y + 1.70, spec["w"], 0.26, spec["cat_title"], 1200, False, GRAY, algn="l")
        bar_chart(slide, spec["x"], card_y + 1.98, spec["w"], spec["h"],
                  spec["cats"], spec["vals"], spec.get("hl_idx", -1), spec.get("hl", "a1"),
                  spec.get("arrow"), spec.get("note"))


def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 底 + 装饰三角(用浅底近似)
    _append(slide, _load("bz-light-bg.xml"))
    _textbox(slide, 0.74, 0.5, 8.0, 0.7, "建筑行业体量持续走高", 3200, True, NAVY, algn="l")
    # 卡1
    stat_card(slide, 0.84, "2020年",
              [("建筑业总产值", "26.39万亿", A1), ("同比增长", "6.2%", A1)],
              [dict(cat_title="建筑业总产值", x=0.98, w=3.21, h=2.7, cats=["2018", "2019", "2020"],
                    vals=[23.1, 24.84, 28.39], hl_idx=-1, hl="a1", arrow="up")])
    # 卡2
    stat_card(slide, 4.88, "2020年6月底",
              [("有施工活动的企业", "102712个", A1), ("同比增长", "10.76%", A1)],
              [dict(cat_title="中国建筑业企业数量", x=5.02, w=3.21, h=2.7, cats=["2018H1", "2019H1", "2020H1"],
                    vals=[65993, 92733, 102712], hl_idx=-1, hl="a2", arrow="up")])
    # 卡3：两个迷你图
    stat_card(slide, 8.92, "2020年6月底",
              [("建筑业从业人数", "4120.9万人", NAVY), ("劳动生产率", "增长4.78%", A1)],
              [dict(cat_title="从业人数（万人）", x=9.06, w=1.55, h=2.7, cats=["2019H2", "2020H2"],
                    vals=[43.1, 41.2], hl_idx=None, hl="gray", arrow="down", note="下降4.78%"),
               dict(cat_title="劳动生产率（%）", x=10.72, w=1.55, h=2.7, cats=["2019H2", "2020H2"],
                    vals=[22.1, 23.2], hl_idx=None, hl="gray")])
    out = os.path.join(ROOT, "_gallery", "data-chart.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    make_demo()
