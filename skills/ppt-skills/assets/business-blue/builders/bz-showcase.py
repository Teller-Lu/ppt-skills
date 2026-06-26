# -*- coding: utf-8 -*-
"""可复用 图文产品 showcase builder（business-blue 风格，蒸自 work_inspire_128 p97）。
版式 = 左两张圆角倒影图槽叠层 + 右(眉标胶囊 + 标题 + 正文 + N 特性 chip) + 左右装饰涡卷群。
背景复用 vectors/bz-light-bg.xml；装饰 = vectors/bz-deco-swirl.xml（本页蒸的一件,左原样 + 右 flipH/flipV 镜像出血）。
图槽/眉标/chip/图标在 builder 原生画。内容（图片/标题/正文/特性）用时传入、不写死（§6.24）。

调用：build_showcase(slide, images, eyebrow, title, body, features)
  images=[main_photo_or_None, overlap_photo_or_None]（两图槽，None=占位）
  features=[(glyph, label), ...]  glyph ∈ {robot, chart, bulb}
直接运行 = p97 原内容 demo → _gallery/showcase.pptx。"""
import os
import sys
import copy
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
NAVY = "0A2A55"
A1 = "0165FF"
A2 = "00BEFF"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
VEC = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue", "vectors")
_parser = etree.XMLParser(resolve_entities=False, no_network=True)
_uid = [1000]


def _nid():
    _uid[0] += 1
    return _uid[0]


def _load(name):
    return etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)


def _append(slide, el):
    slide.shapes._spTree.append(el)
    return el


def _emu(v):
    return int(round(v * IN))


# ---------- 背景：浅色底 + 涡卷装饰群（同一件，左原样 / 右镜像 flipH+flipV 缩放出血）----------
_SWIRL_L = ((112355, 286417), (4580479, 5532484), False, False)
_SWIRL_R = ((10940507, 3985342), (2502985, 3023204), True, True)


def _swirl(slide, spec):
    (ox, oy), (cx, cy), fh, fv = spec
    el = _load("bz-deco-swirl.xml")
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    xf = el.find(f"{P}grpSpPr/{A}xfrm")
    xf.find(f"{A}off").set("x", str(ox)); xf.find(f"{A}off").set("y", str(oy))
    xf.find(f"{A}ext").set("cx", str(cx)); xf.find(f"{A}ext").set("cy", str(cy))
    if fh:
        xf.set("flipH", "1")
    if fv:
        xf.set("flipV", "1")
    _append(slide, el)


def add_bg(slide):
    bg = _load("bz-light-bg.xml")
    bg.find(f".//{P}cNvPr").set("id", str(_nid()))
    _append(slide, bg)
    _swirl(slide, _SWIRL_L)
    _swirl(slide, _SWIRL_R)


# ---------- 圆角倒影 图槽（占位 azure 渐变 + 图片图标；或 blipFill 填真图）----------
def image_slot(slide, x, y, w, h, adj=3400, photo=None):
    refl = ('<a:effectLst><a:reflection blurRad="190500" stA="42000" stPos="0" '
            'endA="300" endPos="52000" dist="0" dir="5400000" sy="-100000" '
            'algn="bl" rotWithShape="0"/></a:effectLst>')
    if photo:
        img_part, rId = slide.part.get_or_add_image_part(photo)
        fill = f'<p:blipFill {nsdecls("r")}><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        sp = parse_xml(
            f'<p:pic {nsdecls("p", "a", "r")}><p:nvPicPr><p:cNvPr id="{_nid()}" name="img-slot"/>'
            f'<p:cNvPicPr/><p:nvPr/></p:nvPicPr>{fill}'
            f'<p:spPr><a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
            f'<a:ln w="19050"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>{refl}</p:spPr></p:pic>')
        return _append(slide, sp)
    # 占位：azure 渐变 roundRect + 白边 + 倒影
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="img-slot"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A2}"><a:lumMod val="60000"/><a:lumOff val="20000"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{A1}"><a:lumMod val="80000"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="2700000" scaled="0"/></a:gradFill>'
        f'<a:ln w="19050"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>{refl}</p:spPr></p:sp>')
    _append(slide, sp)
    # 白色「图片」图标（山+日）居中
    cx, cy = x + w / 2, y + h / 2
    gw = 0.9
    glyph = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="img-glyph-frame"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(cx - gw / 2)}" y="{_emu(cy - gw * 0.36)}"/><a:ext cx="{_emu(gw)}" cy="{_emu(gw * 0.72)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 8000"/></a:avLst></a:prstGeom>'
        f'<a:noFill/><a:ln w="19050"><a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="80000"/></a:srgbClr></a:solidFill></a:ln></p:spPr></p:sp>')
    _append(slide, glyph)
    sun = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="img-glyph-sun"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(cx - gw * 0.30)}" y="{_emu(cy - gw * 0.22)}"/><a:ext cx="{_emu(gw * 0.16)}" cy="{_emu(gw * 0.16)}"/></a:xfrm>'
        f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="80000"/></a:srgbClr></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    _append(slide, sun)
    mount = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="img-glyph-mount"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(cx - gw * 0.40)}" y="{_emu(cy - gw * 0.02)}"/><a:ext cx="{_emu(gw * 0.80)}" cy="{_emu(gw * 0.34)}"/></a:xfrm>'
        f'<a:prstGeom prst="triangle"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="80000"/></a:srgbClr></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    _append(slide, mount)


# ---------- 眉标渐变胶囊 ----------
def eyebrow(slide, x, y, w, h, text):
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="eyebrow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 50000"/></a:avLst></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A1}"/></a:gs><a:gs pos="100000"><a:srgbClr val="{A2}"/></a:gs>'
        f'</a:gsLst><a:lin ang="0" scaled="0"/></a:gradFill><a:ln><a:noFill/></a:ln>'
        f'<a:effectLst><a:outerShdw blurRad="127000" dist="38100" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="{A1}"><a:lumMod val="75000"/><a:alpha val="35000"/></a:srgbClr></a:outerShdw></a:effectLst></p:spPr>'
        f'<p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="zh-CN" sz="1400" b="1" dirty="0"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>')
    return _append(slide, sp)


# ---------- 文本框 ----------
def textbox(slide, x, y, w, h, text, sz, bold, color, anchor="t"):
    battr = ' b="1"' if bold else ''
    paras = ""
    for i, line in enumerate(text.split("\n")):
        paras += (f'<a:p><a:pPr><a:lnSpc><a:spcPct val="120000"/></a:lnSpc></a:pPr><a:r><a:rPr lang="zh-CN" sz="{sz}"{battr} dirty="0">'
                  f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                  f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
                  f'<a:t>{escape(line)}</a:t></a:r></a:p>')
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="txt"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" rtlCol="0" anchor="{anchor}"><a:normAutofit/></a:bodyPr>{paras}</p:txBody></p:sp>')
    return _append(slide, sp)


# ---------- 特性图标（azure 渐变小方块 + 白色简glyph）----------
def _icon_tile(slide, cx, cy, glyph):
    s = 0.46
    tile = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="feat-icon"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(cx - s / 2)}" y="{_emu(cy - s / 2)}"/><a:ext cx="{_emu(s)}" cy="{_emu(s)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 30000"/></a:avLst></a:prstGeom>'
        f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A2}"/></a:gs><a:gs pos="100000"><a:srgbClr val="{A1}"/></a:gs>'
        f'</a:gsLst><a:lin ang="2700000" scaled="0"/></a:gradFill><a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    _append(slide, tile)

    def w_rect(dx, dy, ww, hh, prst="rect", ad=None):
        av = f'<a:gd name="adj" fmla="val {ad}"/>' if ad is not None else ''
        return parse_xml(
            f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="g"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{_emu(cx + dx)}" y="{_emu(cy + dy)}"/><a:ext cx="{_emu(ww)}" cy="{_emu(hh)}"/></a:xfrm>'
            f'<a:prstGeom prst="{prst}"><a:avLst>{av}</a:avLst></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    if glyph == "chart":
        for i, (dx, hh) in enumerate(((-0.10, 0.10), (-0.02, 0.16), (0.06, 0.22))):
            _append(slide, w_rect(dx, 0.11 - hh, 0.05, hh, ad=20000))
    elif glyph == "robot":
        _append(slide, w_rect(-0.11, -0.06, 0.22, 0.16, "roundRect", 25000))   # 头
        _append(slide, w_rect(-0.055, -0.02, 0.03, 0.03, "ellipse"))           # 眼
        _append(slide, w_rect(0.025, -0.02, 0.03, 0.03, "ellipse"))            # 眼
        _append(slide, w_rect(-0.01, -0.13, 0.02, 0.05))                        # 天线
    else:  # bulb/dot
        _append(slide, w_rect(-0.05, -0.05, 0.10, 0.10, "ellipse"))


def feature_chip(slide, x, y, w, h, glyph, label):
    chip = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="feat-chip"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 18000"/></a:avLst></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="55000"/></a:srgbClr></a:solidFill>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="{A1}"><a:alpha val="45000"/></a:srgbClr></a:solidFill></a:ln></p:spPr></p:sp>')
    _append(slide, chip)
    _icon_tile(slide, x + 0.42, y + h / 2, glyph)
    textbox(slide, x + 0.72, y, w - 0.8, h, label, 1600, True, NAVY, anchor="ctr")


# ---------- 主入口 ----------
def build_showcase(slide, images, eyebrow_text, title, body, features):
    image_slot(slide, 2.71, 1.28, 3.70, 4.33, adj=3311, photo=images[0] if len(images) > 0 else None)
    image_slot(slide, 0.87, 2.68, 3.45, 4.14, adj=3547, photo=images[1] if len(images) > 1 else None)
    eyebrow(slide, 7.36, 1.58, 1.30, 0.40, eyebrow_text)
    textbox(slide, 7.28, 2.05, 5.14, 1.10, title, 2600, True, NAVY)   # 图文主文案 26（非页标题→§1.3 例外;26 防压正文）
    textbox(slide, 7.24, 3.66, 5.07, 1.55, body, 1400, False, "44546A")   # 下移避让 3 行主文案(防压字)
    xs = [7.24, 9.60]
    ws = [2.20, 2.20]
    for i, (g, lab) in enumerate(features[:2]):
        feature_chip(slide, xs[i], 5.47, ws[i], 0.78, g, lab)


def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    build_showcase(
        slide,
        images=[None, None],
        eyebrow_text="人工智能",
        title="在商业职场中的应用已成为提升工作效率和决策质量的关键因素",
        body="通过自动化重复性高的任务，如数据录入和报告生成，显著提高了工作效率，让员工得以专注于更具创造性和战略性的工作。同时，AI 的数据分析能力为企业决策提供了有力支撑。",
        features=[("robot", "全天候智能助手"), ("chart", "预测分析的前沿应用")],
    )
    out = os.path.join(ROOT, "_gallery", "showcase.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    make_demo()
