# -*- coding: utf-8 -*-
"""可复用 结尾/致谢 closing builder（business-blue，蒸自 work_inspire_128 p128）。
版式 = 浅底 + 城市楼群(四栋,透明,duotone) + 中心点阵球 + 漂浮 azure 球(复用 sphere-glow+虚影) + 大渐变标题 + 英文副标。

蒸的素材(都在 assets/business-blue 里、可单独调用)：
  · rasters/bz-towers.png      四栋玻璃楼群(透明抠黑底)；整体调用 add_towers，或 place_tower(idx) 单栋
  · vectors/bz-dot-sphere.xml  中心点阵球(5655轮廓 custGeom + 底圆，已 recolor srgb)
  · vectors/sphere-glow.xml    立体发光球(复用旧件，本页每球=它+虚影)
背景复用 bz-light-bg。标题/副标属内容、用时传(§6.24)。

title/subtitle 均 wrap=none 单行签名页，按 design-language §1.3 不套 lnSpc（非遗漏）。
调用：build_closing(slide, title="感谢观看", subtitle="THANKS FOR WATCHING!")
直接运行 = demo → _gallery/closing.pptx。"""
import os
import sys
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
A1 = "0165FF"
A2 = "00BEFF"
GRAY = "8A93A3"
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
LIB = os.path.join(ROOT, "skills", "ppt-skills", "assets", "business-blue")
VEC = os.path.join(LIB, "vectors")
TOWERS_PNG = os.path.join(LIB, "rasters", "bz-towers.png")
_parser = etree.XMLParser(resolve_entities=False, no_network=True, huge_tree=True)
_uid = [3000]

# FLOOR §6.18：schemeClr→srgb（保留 lumMod/lumOff/alpha/shade/satMod）
FLOOR = {
    "accent1": "0165FF", "accent2": "00BEFF",
    "bg1": "FFFFFF", "lt1": "FFFFFF", "bg2": "FFFFFF", "lt2": "FFFFFF",
    "dk1": "0A2A55", "tx1": "0A2A55", "dk2": "0A2A55", "tx2": "0A2A55",
    "phClr": "0165FF",
}

# 四栋楼群在透明整图里的 bbox 分数(split_towers.py 算出，供 srcRect 单栋调用)
TOWERS = [(0.0, 0.1903, 0.263, 0.6936), (0.752, 0.2202, 1.0, 0.7086),
          (0.0, 0.5648, 0.448, 1.0), (0.562, 0.5648, 1.0, 1.0)]


def _nid():
    _uid[0] += 1
    return _uid[0]


def _recolor(el):
    for sc in list(el.iter(f"{A}schemeClr")):
        hexv = FLOOR.get(sc.get("val"))
        if hexv:
            sc.tag = f"{A}srgbClr"
            sc.set("val", hexv)


def _load(name, recolor=False):
    el = etree.fromstring(open(os.path.join(VEC, name), "rb").read(), _parser)
    if recolor:
        _recolor(el)
    return el


def _ap(slide, el):
    slide.shapes._spTree.append(el)
    return el


def _e(v):
    return int(round(v * IN))


def add_bg(slide):
    el = _load("bz-light-bg.xml")
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    _ap(slide, el)


# ---------- 楼群(透明四栋) ----------
def _duotone(blip):
    blip.append(parse_xml(
        f'<a:duotone {nsdecls("a")}>'
        f'<a:srgbClr val="{A1}"><a:shade val="55000"/><a:satMod val="120000"/></a:srgbClr>'
        f'<a:srgbClr val="FFFFFF"/></a:duotone>'))
    blip.append(parse_xml(f'<a:alphaModFix {nsdecls("a")} amt="62000"/>'))


def add_towers(slide):
    """四栋一起铺下半幅(透明 + duotone azure 氛围)。"""
    pic = slide.shapes.add_picture(TOWERS_PNG, _e(0), _e(2.28), _e(13.33), _e(5.22))
    _duotone(pic._element.find(f".//{A}blip"))
    return pic


def place_tower(slide, idx, x, y, w, h, duo=True):
    """单栋调用:idx∈0..3(左上/右上/中下左/中下右),裁出该栋放到 (x,y,w,h)。"""
    fx0, fy0, fx1, fy1 = TOWERS[idx]
    pic = slide.shapes.add_picture(TOWERS_PNG, _e(x), _e(y), _e(w), _e(h))
    pic.crop_left = fx0; pic.crop_top = fy0
    pic.crop_right = 1 - fx1; pic.crop_bottom = 1 - fy1
    if duo:
        _duotone(pic._element.find(f".//{A}blip"))
    return pic


# ---------- 中心点阵球 + 柔光底 ----------
def add_dotsphere(slide):
    # 大柔光晕垫底(景深)
    cx, cy = 6.67, 3.55
    glow = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="halo-glow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(cx-3.5)}" y="{_e(cy-3.5)}"/><a:ext cx="{_e(7.0)}" cy="{_e(7.0)}"/></a:xfrm>'
        f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'<a:gradFill><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A2}"><a:alpha val="20000"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{A2}"><a:alpha val="0"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:path path="circle"><a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr></p:sp>')
    _ap(slide, glow)
    # 点阵球(已 recolor srgb，自带绝对 off 居中)
    el = _load("bz-dot-sphere.xml")
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    _ap(slide, el)


# ---------- 漂浮球 = sphere-glow(recolor) + 虚影 ----------
def _sphere(slide, x, y, d):
    # 虚影(更大更淡、柔边，垫在球后)
    gd = d * 1.4
    ghost = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="sphere-ghost"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x-(gd-d)/2)}" y="{_e(y-(gd-d)/2+d*0.06)}"/><a:ext cx="{_e(gd)}" cy="{_e(gd)}"/></a:xfrm>'
        f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'<a:gradFill><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A2}"><a:alpha val="26000"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{A1}"><a:alpha val="0"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:path path="circle"><a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path></a:gradFill>'
        f'<a:ln><a:noFill/></a:ln><a:effectLst><a:softEdge rad="120000"/></a:effectLst></p:spPr></p:sp>')
    _ap(slide, ghost)
    # 立体发光球(复用 sphere-glow.xml，recolor srgb，缩放到 d)
    el = _load("sphere-glow.xml", recolor=True)
    el.find(f".//{P}cNvPr").set("id", str(_nid()))
    xf = el.find(f"{P}spPr/{A}xfrm")
    xf.find(f"{A}off").set("x", str(_e(x))); xf.find(f"{A}off").set("y", str(_e(y)))
    xf.find(f"{A}ext").set("cx", str(_e(d))); xf.find(f"{A}ext").set("cy", str(_e(d)))
    _ap(slide, el)


def add_spheres(slide):
    _sphere(slide, 11.2, -0.38, 3.48)   # 右上大球(出血)
    _sphere(slide, 0.59, 0.45, 1.64)    # 左上球
    _sphere(slide, 1.5, 5.42, 0.72)     # 左下小球
    _sphere(slide, 11.72, 5.08, 0.92)   # 右下小球


def title(slide, text, sz=13800):
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="title"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(1.0)}" y="{_e(1.55)}"/><a:ext cx="{_e(11.33)}" cy="{_e(2.4)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="none" anchor="ctr"/><a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="zh-CN" sz="{sz}" b="1" dirty="0">'
        f'<a:gradFill><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{A1}"/></a:gs><a:gs pos="100000"><a:srgbClr val="{A2}"/></a:gs>'
        f'</a:gsLst><a:lin ang="3000000" scaled="1"/></a:gradFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>')
    _ap(slide, sp)


def subtitle(slide, text):
    sp = parse_xml(
        f'<p:sp {nsdecls("p", "a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="subtitle"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(1.0)}" y="{_e(3.95)}"/><a:ext cx="{_e(11.33)}" cy="{_e(0.5)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="none" anchor="ctr"/><a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="en-US" sz="1600" b="0" spc="600" dirty="0">'
        f'<a:solidFill><a:srgbClr val="{GRAY}"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>')
    _ap(slide, sp)


def build_closing(slide, title_text="感谢观看", subtitle_text="THANKS FOR WATCHING!", title_sz=13800):
    add_bg(slide)
    add_towers(slide)
    add_dotsphere(slide)
    add_spheres(slide)
    title(slide, title_text, title_sz)
    subtitle(slide, subtitle_text)


def make_demo():
    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_closing(slide, "感谢观看", "THANKS FOR WATCHING!")
    out = os.path.join(ROOT, "_gallery", "closing.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    make_demo()
