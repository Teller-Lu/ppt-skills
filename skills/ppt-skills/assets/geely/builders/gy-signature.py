# -*- coding: utf-8 -*-
"""GEELY 公司风 · 签名页 builder（蒸自 四驱扭矩寻优分配系统专项汇报.pptx + 驱动型式对弯道性能的仿真分析_E22H.pptx）。
一个 builder 4 版式：
  build_cover_white  纯白极简封面（白底 + 左上联名logo占位 + 竖线引导 + 黑标题40pt左对齐 + 灰日期15pt）
  build_cover_silk   深灰蓝丝绸封面（原生重建丝绸折面背景 + 右上logo占位 + 白标题36pt粗居中 + 部门/编制/时间）
  build_closing      THANKS 结尾（白底 + 居中双logo占位 + THANKS 灰70pt 居中）
  build_content      中间页角标框架（白底 + 左上标题配浅蓝强调块 + 右上logo占位 + 空内容区，左右上各一角标）

纪律（蒸馏流程 §0 步4 + 记忆）：
  · 微软雅黑烤进每个 run 的 latin/ea/cs（ea 管中文必设，防注入回退宋体）
  · DLP 水印（EagleCloudWatermark / 保密:一级）一律不蒸
  · 真实 logo 不入库 —— 全部留虚线占位框，实战时换用户真实 logo
  · 丝绸折面背景是可复用母题，调好后固化入 vectors/gy-silk-bg.xml（当前 inline 调试）
demo: python gy-signature.py [white|silk|closing|content]
"""
import os, sys
from xml.sax.saxutils import escape
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
sys.stdout.reconfigure(encoding="utf-8")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
IN = 914400
FACE = "微软雅黑"
INK = "262626"       # 主标题近黑
GRAY = "595959"      # 副文字灰（日期/部门），实测自 p1
THX = "4D4D4D"       # THANKS 中深灰
LBLUE = "BDD7EE"     # 中间页浅蓝强调块
VLINE = "404040"     # 竖线引导深灰
PLACE = "BFBFBF"     # 占位框线灰
PLACE_TX = "A6A6A6"  # 占位提示字灰
WHITE = "FFFFFF"

CW, CH = 13.33, 7.5
ROOT = r"D:\Lu.Yao7\VehicleTellerLu\Agent\Shared\skills-creator\ppt-skills"
LIB = os.path.join(ROOT, "skills", "ppt-skills", "assets", "geely")
VEC = os.path.join(LIB, "vectors")
RAS = os.path.join(LIB, "rasters")
SILK_BG = os.path.join(RAS, "gy-cover-silk-bg.jpg")  # 原模板丝绸背景图(原样保留,不改动)
_parser = etree.XMLParser(resolve_entities=False, no_network=True, huge_tree=True)
_uid = [5000]


def _nid():
    _uid[0] += 1
    return _uid[0]


def _e(v):
    return int(round(v * IN))


def _ap(slide, el):
    slide.shapes._spTree.append(el)
    return el


# ---------- 文字（烤微软雅黑 ea） ----------
def _text(slide, x, y, w, h, text, sz, color, bold=False, algn="l", anchor="t", spc=None):
    spc_attr = f' spc="{spc}"' if spc else ""
    sp = parse_xml(
        f'<p:sp {nsdecls("p","a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="tx"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" anchor="{anchor}"/><a:p><a:pPr algn="{algn}"/>'
        f'<a:r><a:rPr lang="zh-CN" sz="{sz}" b="{1 if bold else 0}"{spc_attr} dirty="0">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>')
    return _ap(slide, sp)


def _text_lines(slide, x, y, w, h, lines, sz, color, bold=False, algn="l", anchor="t"):
    ps = ""
    for ln in lines:
        ps += (f'<a:p><a:pPr algn="{algn}"/><a:r><a:rPr lang="zh-CN" sz="{sz}" b="{1 if bold else 0}" dirty="0">'
               f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
               f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
               f'<a:t>{escape(ln)}</a:t></a:r></a:p>')
    sp = parse_xml(
        f'<p:sp {nsdecls("p","a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="tx-multi"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" anchor="{anchor}"/>{ps}</p:txBody></p:sp>')
    return _ap(slide, sp)


# ---------- 背景 ----------
def add_white_bg(slide):
    sp = parse_xml(
        f'<p:sp {nsdecls("p","a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="bg-white"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{_e(CW)}" cy="{_e(CH)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{WHITE}"/></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:p/></p:txBody></p:sp>')
    return _ap(slide, sp)


def add_silk_photo_bg(slide):
    """丝绸封面背景 = 原模板背景图，原样铺满（不改动；右上自带 GEELY 字标）。"""
    slide.shapes.add_picture(SILK_BG, 0, 0, _e(CW), _e(CH))


# ---------- 装饰件 ----------
def logo_placeholder(slide, x, y, w, h, label="LOGO 占位", line=PLACE, tx=PLACE_TX):
    sp = parse_xml(
        f'<p:sp {nsdecls("p","a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="logo-ph"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="{line}"/></a:solidFill><a:prstDash val="dash"/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" anchor="ctr"/><a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="zh-CN" sz="1000" dirty="0"><a:solidFill><a:srgbClr val="{tx}"/></a:solidFill>'
        f'<a:latin typeface="{FACE}"/><a:ea typeface="{FACE}"/><a:cs typeface="{FACE}"/></a:rPr>'
        f'<a:t>{escape(label)}</a:t></a:r></a:p></p:txBody></p:sp>')
    return _ap(slide, sp)


def vline(slide, x, y, h, color=VLINE, w_emu=12700):
    sp = parse_xml(
        f'<p:cxnSp {nsdecls("p","a")}><p:nvCxnSpPr><p:cNvPr id="{_nid()}" name="vline"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="0" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{w_emu}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln></p:spPr></p:cxnSp>')
    return _ap(slide, sp)


def hline(slide, x, y, w, color=VLINE, w_emu=12700):
    sp = parse_xml(
        f'<p:cxnSp {nsdecls("p","a")}><p:nvCxnSpPr><p:cNvPr id="{_nid()}" name="hline"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{w_emu}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln></p:spPr></p:cxnSp>')
    return _ap(slide, sp)


def accent_block(slide, x, y, w, h, color=LBLUE):
    sp = parse_xml(
        f'<p:sp {nsdecls("p","a")}><p:nvSpPr><p:cNvPr id="{_nid()}" name="accent"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{_e(x)}" y="{_e(y)}"/><a:ext cx="{_e(w)}" cy="{_e(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill><a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:p/></p:txBody></p:sp>')
    return _ap(slide, sp)


# ---------- 4 版式 ----------
def build_cover_white(slide, title, date_text):
    add_white_bg(slide)
    logo_placeholder(slide, 1.25, 1.2, 4.3, 1.45, "联名 LOGO 占位（用时换真实 logo）")
    vline(slide, 1.55, 3.95, 1.78, color=VLINE, w_emu=15875)
    _text(slide, 1.85, 3.95, 10.2, 1.0, title, 4000, INK, bold=False, algn="l", anchor="ctr")
    _text(slide, 1.88, 5.02, 6.0, 0.45, date_text, 1500, GRAY, algn="l", anchor="ctr")


def build_cover_silk(slide, title, dept, author, date_text):
    add_silk_photo_bg(slide)   # 原模板背景图自带右上 GEELY，无需再放 logo 占位
    # 文字严格复刻原模板：标题居中 + 标题左侧竖线引导 + 副标题左对齐
    _text(slide, 1.39, 3.92, 9.98, 0.78, title, 3600, WHITE, bold=True, algn="ctr", anchor="ctr")
    vline(slide, 1.2, 3.92, 1.62, color="D5DAE0", w_emu=15875)
    _text_lines(slide, 1.41, 4.72, 7.06, 1.3,
                [f"部门：{dept}", f"编制：{author}", f"时间：{date_text}"],
                1600, WHITE, algn="l", anchor="t")


def build_closing(slide, thanks="THANKS"):
    add_white_bg(slide)
    logo_placeholder(slide, 4.77, 2.08, 3.79, 1.46, "双 LOGO 占位（用时换真实 logo）")
    _text(slide, 0.0, 3.85, 13.33, 1.35, thanks, 7000, THX, bold=False, algn="ctr", anchor="ctr", spc=300)


def build_content(slide, page_title="技术图谱"):
    add_white_bg(slide)
    accent_block(slide, 0.0, 0.44, 0.92, 0.42, LBLUE)
    _text(slide, 0.46, 0.12, 10.6, 0.92, page_title, 2800, INK, bold=False, algn="l", anchor="ctr")
    logo_placeholder(slide, 11.38, 0.3, 1.55, 0.5, "LOGO")
    # 内容区留空（用户填）


# ---------- demo ----------
def _demo(mode):
    prs = Presentation()
    prs.slide_width = Emu(_e(CW))
    prs.slide_height = Emu(_e(CH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if mode == "white":
        build_cover_white(slide, "四驱扭矩寻优分配系统专项汇报", "2026 年 06 月 23 日")
    elif mode == "silk":
        build_cover_silk(slide, "驱动型式对弯道性能的仿真分析_E22H", "系统工程部", "李林儒", "2026 年 5 月 18 日")
    elif mode == "closing":
        build_closing(slide, "THANKS")
    else:
        build_content(slide, "技术图谱")
    out = os.path.join(ROOT, "_experiments", "company-style", f"gy_{mode}.pptx")
    prs.save(out)
    print("saved", out)


def make_master():
    """装配 4 页 geely 合集册（封面白 → 封面丝绸 → 中间页 → 结尾）→ _gallery/geely-master.pptx。"""
    prs = Presentation()
    prs.slide_width = Emu(_e(CW))
    prs.slide_height = Emu(_e(CH))
    blank = prs.slide_layouts[6]
    build_cover_white(prs.slides.add_slide(blank), "四驱扭矩寻优分配系统专项汇报", "2026 年 06 月 23 日")
    build_cover_silk(prs.slides.add_slide(blank), "驱动型式对弯道性能的仿真分析_E22H", "系统工程部", "李林儒", "2026 年 5 月 18 日")
    build_content(prs.slides.add_slide(blank), "技术图谱")
    build_closing(prs.slides.add_slide(blank), "THANKS")
    out = os.path.join(ROOT, "_gallery", "geely-master.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    arg = sys.argv[1] if len(sys.argv) > 1 else "white"
    if arg == "master":
        make_master()
    else:
        _demo(arg)
